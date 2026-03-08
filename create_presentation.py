"""
Generate a professional PPTX presentation for the NIELIT QR Attendance System.
Presenter: Tanya Singh | NIELIT Intern
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os

# ── Constants ────────────────────────────────────────────────
NAVY      = RGBColor(0x00, 0x2B, 0x5C)   # Deep navy
DARK_BLUE = RGBColor(0x00, 0x66, 0xB3)   # NIELIT blue
LIGHT_BLUE= RGBColor(0x00, 0xA0, 0xE3)   # Accent cyan
GREEN     = RGBColor(0x8D, 0xC6, 0x3F)   # Accent green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY  = RGBColor(0x66, 0x66, 0x66)
PINK      = RGBColor(0xE9, 0x1E, 0x8C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

HEADING_FONT = "Segoe UI"
BODY_FONT    = "Segoe UI"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

today_str = datetime.date.today().strftime("%B %d, %Y")   # March 08, 2026

# ── Helper functions ─────────────────────────────────────────

def add_solid_bg(slide, color):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=BODY_FONT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_bullet_slide_content(tf, items, font_size=16, color=DARK_GRAY, bold_prefix=True):
    """Append bullet items to an existing text frame. Items can be str or (bold_part, rest)."""
    for item in items:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(4)
        p.line_spacing = Pt(int(font_size * 1.5))
        if isinstance(item, tuple):
            run_b = p.add_run()
            run_b.text = item[0]
            run_b.font.size = Pt(font_size)
            run_b.font.bold = True
            run_b.font.color.rgb = DARK_BLUE
            run_b.font.name = BODY_FONT
            run_n = p.add_run()
            run_n.text = item[1]
            run_n.font.size = Pt(font_size)
            run_n.font.color.rgb = color
            run_n.font.name = BODY_FONT
        else:
            run = p.add_run()
            run.text = f"  \u2022  {item}"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = BODY_FONT


def add_slide_number(slide, num):
    add_textbox(slide, SLIDE_W - Inches(1), SLIDE_H - Inches(0.5),
                Inches(0.8), Inches(0.4), str(num),
                font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def make_section_slide(title_text, subtitle_text=None, slide_num=1):
    """Navy background section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_solid_bg(slide, NAVY)

    # Accent bar
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), GREEN)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                title_text, font_size=40, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=HEADING_FONT)
    if subtitle_text:
        add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                    subtitle_text, font_size=20, color=LIGHT_BLUE,
                    alignment=PP_ALIGN.CENTER)
    add_slide_number(slide, slide_num)
    return slide


def make_content_slide(title_text, slide_num):
    """White background slide with navy header bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)

    # Top header bar
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.8),
                title_text, font_size=28, color=WHITE, bold=True,
                font_name=HEADING_FONT)

    # Green accent under header
    add_shape_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.05), GREEN)

    # Bottom bar
    add_shape_rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), NAVY)
    add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.33), Inches(5), Inches(0.3),
                "NIELIT QR Attendance System  |  Tanya Singh",
                font_size=9, color=LIGHT_BLUE)
    add_slide_number(slide, slide_num)
    return slide


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

# Large accent rectangle on left
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.25), SLIDE_H, GREEN)

# Decorative circle
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(5), Inches(2), Inches(2))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0x00, 0x44, 0x88)
shape2.line.fill.background()

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(9), Inches(1.2),
            "QR Attendance System", font_size=48, color=WHITE, bold=True,
            font_name=HEADING_FONT)

# Subtitle
add_textbox(slide, Inches(1), Inches(2.8), Inches(9), Inches(0.6),
            "Dynamic QR-Based Student Attendance Management",
            font_size=22, color=LIGHT_BLUE, font_name=HEADING_FONT)

# Green accent line
add_shape_rect(slide, Inches(1), Inches(3.6), Inches(3), Inches(0.06), GREEN)

# Presented by
add_textbox(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.5),
            "NIELIT Internship Presentation", font_size=18, color=OFF_WHITE)

add_textbox(slide, Inches(1), Inches(4.7), Inches(8), Inches(0.5),
            "Presented by:  Tanya Singh  |  NIELIT Intern",
            font_size=20, color=GREEN, bold=True)

add_textbox(slide, Inches(1), Inches(5.4), Inches(8), Inches(0.5),
            f"Date: {today_str}", font_size=16, color=MID_GRAY)

add_slide_number(slide, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction / Overview
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Introduction & Overview", 2)

# Problem statement box
prob_box = add_shape_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5), OFF_WHITE)
txBox = add_textbox(slide, Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.4),
                    "\U0001F3AF  Problem Statement", font_size=20, color=DARK_BLUE, bold=True)
txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5.4), Inches(1.8))
tf = txBox2.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, [
    "Manual attendance is time-consuming and error-prone",
    "Proxy attendance is easy with paper-based systems",
    "No real-time tracking or exportable digital records",
    "Teachers lack tools for monthly analytics & reports",
], font_size=14)

# Objectives box
obj_box = add_shape_rect(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), OFF_WHITE)
add_textbox(slide, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.4),
            "\U0001F680  Objectives", font_size=20, color=DARK_BLUE, bold=True)
txBox3 = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.4), Inches(1.8))
tf2 = txBox3.text_frame
tf2.word_wrap = True
add_bullet_slide_content(tf2, [
    "Build a secure QR-based attendance system",
    "15-second rotating QR to prevent proxy attendance",
    "Separate portals for Students, Teachers & Admins",
    "Export monthly reports as PDF & CSV",
], font_size=14)

# Overview paragraph
add_textbox(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(2.5),
            "The QR Attendance System is a full-stack web application that replaces traditional roll-call "
            "with dynamic, time-validated QR codes. Teachers generate a rotating QR code on-screen; students "
            "scan it using their phone camera within a 15-second window. The system records attendance with "
            "exact scan timestamps, prevents duplicate marking, and provides real-time dashboards with "
            "monthly PDF/CSV export capabilities. Deployed on Netlify (frontend) and Supabase (backend).",
            font_size=15, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Technology Stack
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Technology Stack & Tools", 3)

tech_categories = [
    ("Frontend", [
        ("React 18", "Core UI library"),
        ("Vite 5", "Build tool & dev server"),
        ("React Router v6", "Client-side routing"),
        ("html5-qrcode", "QR scanning via camera"),
        ("qrcode", "QR code generation"),
        ("jsPDF", "PDF report generation"),
    ], DARK_BLUE),
    ("Backend", [
        ("Supabase Edge Functions", "Serverless API (Deno/TS)"),
        ("Supabase Auth", "JWT authentication"),
        ("PostgreSQL", "Relational database"),
        ("Row Level Security", "Database access control"),
    ], RGBColor(0x00, 0x80, 0x60)),
    ("DevOps & Tools", [
        ("Netlify", "Frontend hosting & CI/CD"),
        ("Git & GitHub", "Version control"),
        ("VS Code", "Development IDE"),
        ("Supabase Dashboard", "DB management & monitoring"),
    ], RGBColor(0x8B, 0x5C, 0xF6)),
]

col_x_positions = [Inches(0.5), Inches(4.6), Inches(8.7)]
for idx, (cat_name, items, accent_color) in enumerate(tech_categories):
    x = col_x_positions[idx]
    # Card background
    add_shape_rect(slide, x, Inches(1.4), Inches(3.8), Inches(5.4), OFF_WHITE)
    # Accent top bar on card
    add_shape_rect(slide, x, Inches(1.4), Inches(3.8), Inches(0.08), accent_color)

    add_textbox(slide, x + Inches(0.2), Inches(1.55), Inches(3.4), Inches(0.5),
                cat_name, font_size=20, color=accent_color, bold=True)

    y_start = Inches(2.2)
    for i, (tech, desc) in enumerate(items):
        y = y_start + Inches(i * 0.72)
        add_textbox(slide, x + Inches(0.3), y, Inches(3.2), Inches(0.3),
                    tech, font_size=14, color=DARK_GRAY, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.28), Inches(3.2), Inches(0.3),
                    desc, font_size=11, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("System Architecture", 4)

# ── Three-tier architecture diagram built with shapes ──
tiers = [
    ("Frontend  (React + Vite)", "Student Portal  |  Teacher Portal  |  Admin Panel\nHosted on Netlify", DARK_BLUE, Inches(1.5)),
    ("Backend  (Supabase Edge Functions)", "attendance-api  |  manage-users  |  mark-attendance\nServerless Deno/TypeScript", RGBColor(0x00, 0x80, 0x60), Inches(3.3)),
    ("Database  (PostgreSQL)", "students  |  courses  |  attendance  |  profiles\nRow Level Security  |  Supabase Auth", RGBColor(0x8B, 0x5C, 0xF6), Inches(5.1)),
]

for (title, desc, color, y) in tiers:
    box = add_shape_rect(slide, Inches(2), y, Inches(9.3), Inches(1.4), color)
    add_textbox(slide, Inches(2.3), y + Inches(0.1), Inches(8.7), Inches(0.5),
                title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.3), y + Inches(0.6), Inches(8.7), Inches(0.7),
                desc, font_size=13, color=OFF_WHITE, alignment=PP_ALIGN.CENTER)

# Arrows between tiers
for y_arrow in [Inches(2.95), Inches(4.75)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), y_arrow, Inches(0.5), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()

# Labels on left
add_textbox(slide, Inches(0.3), Inches(1.7), Inches(1.5), Inches(0.5),
            "Presentation\nLayer", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.3), Inches(3.5), Inches(1.5), Inches(0.5),
            "Application\nLayer", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.3), Inches(5.3), Inches(1.5), Inches(0.5),
            "Data\nLayer", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features (Part 1)
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Key Features", 5)

features = [
    ("\U0001F4F1  Dynamic QR Codes",
     "Auto-regenerating every 15 seconds\nfor maximum security against\nproxy attendance"),
    ("\u23F1\uFE0F  Time-Based Validation",
     "Strict 15-second window ensures\nonly live, in-class students\ncan mark attendance"),
    ("\U0001F468\u200D\U0001F3EB  Dual Portal System",
     "Separate interfaces for Students\n(scan QR) and Teachers\n(generate QR + manage)"),
    ("\U0001F4CA  Real-Time Dashboard",
     "Instant attendance counts,\nmonthly analytics with\npercentage calculations"),
    ("\U0001F4C4  PDF & CSV Export",
     "Calendar-style monthly PDF\nwith NIELIT branding +\nCSV data export"),
    ("\U0001F512  Role-Based Access",
     "Admin, Teacher & Student roles\nwith JWT authentication\nand Supabase RLS"),
]

positions = [
    (Inches(0.5), Inches(1.4)),  (Inches(4.55), Inches(1.4)),  (Inches(8.6), Inches(1.4)),
    (Inches(0.5), Inches(4.2)),  (Inches(4.55), Inches(4.2)),  (Inches(8.6), Inches(4.2)),
]

for i, ((feat_title, feat_desc), (x, y)) in enumerate(zip(features, positions)):
    add_shape_rect(slide, x, y, Inches(3.8), Inches(2.4), OFF_WHITE)
    add_shape_rect(slide, x, y, Inches(3.8), Inches(0.06), DARK_BLUE if i < 3 else GREEN)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5),
                feat_title, font_size=16, color=DARK_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.4),
                feat_desc, font_size=13, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — How It Works (Workflow)
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("How It Works — Attendance Workflow", 6)

steps = [
    ("Step 1", "Teacher Login", "Teacher logs into Faculty Portal\nusing email credentials"),
    ("Step 2", "Start QR Session", "Clicks 'Start Attendance Session'\nQR code appears, refreshing\nevery 15 seconds"),
    ("Step 3", "Student Scans", "Student opens Student Portal,\npoints camera at QR code\non teacher's screen"),
    ("Step 4", "Validation", "System validates QR freshness\n(< 15 sec), checks course match,\nprevents duplicate marking"),
    ("Step 5", "Record Created", "Attendance row inserted in DB\nwith status='present' and\nexact scan_time timestamp"),
    ("Step 6", "Reports & Export", "Teacher views daily/monthly\nreports, exports PDF with\ncalendar grid & stats"),
]

for i, (step_label, step_title, step_desc) in enumerate(steps):
    x = Inches(0.4) + Inches(i * 2.12)
    y = Inches(1.6)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_BLUE if i % 2 == 0 else GREEN
    circle.line.fill.background()
    # Number text in circle
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Arrow between steps (except last)
    if i < 5:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                                        x + Inches(1.55), y + Inches(0.15),
                                        Inches(0.55), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LIGHT_BLUE
        arrow.line.fill.background()

    # Title
    add_textbox(slide, x, y + Inches(0.9), Inches(2.0), Inches(0.4),
                step_title, font_size=13, color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Desc
    add_textbox(slide, x, y + Inches(1.35), Inches(2.0), Inches(1.5),
                step_desc, font_size=10, color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER)

# Connector line behind
add_shape_rect(slide, Inches(0.95), Inches(1.93), Inches(11.4), Inches(0.04), OFF_WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Implementation Details
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Implementation Details", 7)

# Left column: Backend
add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5), OFF_WHITE)
add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.06), DARK_BLUE)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.4),
            "\U0001F5A5\uFE0F  Backend — Supabase Edge Functions", font_size=18, color=DARK_BLUE, bold=True)

backend_items = [
    "3 Edge Functions: attendance-api, manage-users, mark-attendance",
    "Written in TypeScript, runs on Deno serverless runtime",
    "attendance-api handles: server-time, mark-attendance,",
    "    mark-self, get-daily, get-monthly, get-courses, finalize",
    "JWT token verification on every request",
    "Service Role Key for admin-level DB operations",
    "Student lookup via auth metadata → roll_number mapping",
    "Attendance ID format: {student_id}-{date} (ensures uniqueness)",
]
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, backend_items, font_size=12, color=DARK_GRAY)

# Right column: Frontend
add_shape_rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5), OFF_WHITE)
add_shape_rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.06), GREEN)
add_textbox(slide, Inches(7.0), Inches(1.55), Inches(5.6), Inches(0.4),
            "\U0001F310  Frontend — React + Vite", font_size=18, color=RGBColor(0x00, 0x80, 0x60), bold=True)

frontend_items = [
    "Single Page Application with React Router v6",
    "Components: Login, StudentPortal, TeacherPortal, AdminPanel",
    "QR generation using 'qrcode' library (student side)",
    "QR scanning using 'html5-qrcode' camera API",
    "PDF generation with jsPDF — calendar-style monthly grid",
    "Responsive design — works on mobile & desktop",
    "Environment variables via Vite (VITE_SUPABASE_URL, etc.)",
    "Deployed on Netlify with automatic CI/CD from GitHub",
]
txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.6), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_bullet_slide_content(tf2, frontend_items, font_size=12, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Database Schema
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Database Schema", 8)

tables = [
    ("students", "id  |  roll_number  |  name\nemail  |  course_code  |  user_id", DARK_BLUE),
    ("courses", "course_code  |  course_name\nteacher_id  |  teacher_name", RGBColor(0x00, 0x80, 0x60)),
    ("attendance", "id  |  student_id  |  date\nstatus  |  scan_time  |  finalized", RGBColor(0xE6, 0x5C, 0x00)),
    ("profiles", "id (UUID)  |  role\n(admin / teacher / student)", RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (tbl_name, tbl_cols, color) in enumerate(tables):
    x = Inches(0.4) + Inches(i * 3.2)
    y = Inches(1.6)
    # Table card
    add_shape_rect(slide, x, y, Inches(2.9), Inches(2.5), color)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.5),
                tbl_name, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Divider inside card
    add_shape_rect(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.5), Inches(0.03), WHITE)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.8), Inches(2.6), Inches(1.5),
                tbl_cols, font_size=12, color=OFF_WHITE, alignment=PP_ALIGN.CENTER)

# Relationships text
add_textbox(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(2.2),
            "Key Relationships:\n"
            "  \u2022  attendance.student_id \u2192 students.id  (tracks per-student daily presence)\n"
            "  \u2022  students.course_code \u2192 courses.course_code  (course enrollment)\n"
            "  \u2022  courses.teacher_id \u2192 profiles.id  (teacher assignment via Supabase Auth UUID)\n"
            "  \u2022  Attendance ID = '{student_id}-{date}' ensures one record per student per day\n"
            "  \u2022  Supabase Auth handles user creation, JWT tokens & metadata storage",
            font_size=13, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Screenshots: Login & Faculty Portal
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Application Screenshots — Portals", 9)

# Since we can't embed actual images, we create placeholder frames
# showing what each screenshot displays

screens = [
    ("Login Page", "NIELIT branded login with\nemail & password fields\n+ Forgot Password link",
     Inches(0.5), Inches(1.4), Inches(3.8), Inches(5.3)),
    ("Faculty Portal — QR Code", "Teacher generates rotating\nQR code every 15 seconds\nShows Present / Absent / Total\ncounts in real-time",
     Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.3)),
    ("Student Portal — Scan", "Student scans QR via phone\ncamera, sees confirmation\n'Attendance Marked!'\nwith exact timestamp",
     Inches(8.9), Inches(1.4), Inches(3.8), Inches(5.3)),
]

for (scr_title, scr_desc, x, y, w, h) in screens:
    # Frame
    frame = add_shape_rect(slide, x, y, w, h, OFF_WHITE)
    frame.line.color.rgb = LIGHT_BLUE
    frame.line.width = Pt(2)

    # Screenshot area placeholder
    add_shape_rect(slide, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(2.2),
                   RGBColor(0xE8, 0xEC, 0xF0))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), w - Inches(0.4), Inches(1),
                "[Screenshot]", font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # Caption below
    add_textbox(slide, x + Inches(0.15), y + h - Inches(2.0), w - Inches(0.3), Inches(0.4),
                scr_title, font_size=15, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), y + h - Inches(1.5), w - Inches(0.3), Inches(1.3),
                scr_desc, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Screenshots: Monthly Report & PDF
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Results & Output — Reports", 10)

screens2 = [
    ("Monthly Report Portal", "Interactive monthly attendance grid\nshowing P / A / Weekend for each day\nOverall %, Session Days, Student count\nStudent-wise summary table",
     Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.3)),
    ("Exported PDF Report", "NIELIT-branded landscape PDF\nCalendar-style grid: 1-31 columns\nP (green) / A (red) / OFF (weekend)\nPresent count column + Session Days\nGenerated via jsPDF",
     Inches(6.9), Inches(1.4), Inches(5.8), Inches(5.3)),
]

for (scr_title, scr_desc, x, y, w, h) in screens2:
    frame = add_shape_rect(slide, x, y, w, h, OFF_WHITE)
    frame.line.color.rgb = LIGHT_BLUE
    frame.line.width = Pt(2)

    add_shape_rect(slide, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(2.5),
                   RGBColor(0xE8, 0xEC, 0xF0))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), Inches(1),
                "[Screenshot]", font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.15), y + h - Inches(2.3), w - Inches(0.3), Inches(0.4),
                scr_title, font_size=16, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), y + h - Inches(1.8), w - Inches(0.3), Inches(1.6),
                scr_desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Challenges & Learnings
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Challenges & Learnings", 11)

# Left: Challenges
add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), OFF_WHITE)
add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.06), PINK)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.4),
            "\u26A0\uFE0F  Challenges Faced", font_size=20, color=PINK, bold=True)

challenges = [
    "QR timing synchronization between teacher & student devices",
    "Mapping Supabase Auth UUIDs to integer student IDs in the DB",
    "Handling timezone differences for scan_time across IST/UTC",
    "PDF generation with jsPDF — manual layout for calendar grid",
    "Responsive camera scanning on mobile browsers (HTTPS requirement)",
    "Data format mismatch between backend API & frontend components",
    "Weekend/holiday detection logic for accurate working day counts",
]
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.15), Inches(5.4), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
add_bullet_slide_content(tf, challenges, font_size=13)

# Right: Learnings
add_shape_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5), OFF_WHITE)
add_shape_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.06), GREEN)
add_textbox(slide, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.4),
            "\U0001F4A1  Key Learnings", font_size=20, color=RGBColor(0x2E, 0x7D, 0x32), bold=True)

learnings = [
    "Supabase as a complete backend platform (Auth + DB + Functions)",
    "Edge Function development with Deno & TypeScript",
    "React SPA architecture with component-based design",
    "QR code generation & scanning using web APIs",
    "JWT-based authentication & role-based access control",
    "PDF generation techniques using jsPDF (without autoTable)",
    "Git workflow, Netlify CI/CD & production deployment",
]
txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(2.15), Inches(5.4), Inches(4.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
add_bullet_slide_content(tf2, learnings, font_size=13)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusion
# ═══════════════════════════════════════════════════════════════
slide = make_content_slide("Conclusion", 12)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
            "The QR Attendance System successfully demonstrates a modern, secure, and efficient approach "
            "to student attendance management. By leveraging dynamic QR codes with strict time validation, "
            "the system eliminates proxy attendance while providing a seamless experience for both students and teachers.",
            font_size=16, color=DARK_GRAY)

# Summary cards
summary_items = [
    ("\u2705", "Secure", "15-second rotating QR\nprevents proxy attendance"),
    ("\u2705", "Real-Time", "Instant attendance marking\nwith scan timestamps"),
    ("\u2705", "Multi-Portal", "Admin + Teacher + Student\nseparate interfaces"),
    ("\u2705", "Exportable", "PDF & CSV monthly reports\nwith NIELIT branding"),
    ("\u2705", "Deployed", "Live on Netlify + Supabase\naccessible anywhere"),
]

for i, (icon, title, desc) in enumerate(summary_items):
    x = Inches(0.5) + Inches(i * 2.5)
    y = Inches(3.2)
    add_shape_rect(slide, x, y, Inches(2.2), Inches(2.3), OFF_WHITE)
    add_shape_rect(slide, x, y, Inches(2.2), Inches(0.06), GREEN)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.2), Inches(2.0), Inches(0.5),
                f"{icon}  {title}", font_size=16, color=DARK_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.8), Inches(2.0), Inches(1.2),
                desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Future scope
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
            "Future Scope:  Biometric integration  |  Geofencing for location verification  |  "
            "Analytics dashboard with charts  |  Push notifications for students  |  "
            "Multi-institution support",
            font_size=13, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, NAVY)

# Decorative elements
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(-0.5), Inches(3), Inches(3))
circle1.fill.solid()
circle1.fill.fore_color.rgb = DARK_BLUE
circle1.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(5.5), Inches(2.5), Inches(2.5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x00, 0x44, 0x88)
circle2.line.fill.background()

# Main heading
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "Thank You!", font_size=60, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER, font_name=HEADING_FONT)

# Green accent line
add_shape_rect(slide, Inches(5), Inches(3.1), Inches(3.3), Inches(0.06), GREEN)

# Subtitle
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
            "Presented by:  Tanya Singh  |  NIELIT Intern",
            font_size=22, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
            "\"Building secure, modern solutions for smarter attendance management\"",
            font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# Contact info
add_textbox(slide, Inches(2), Inches(5.2), Inches(9), Inches(1.2),
            "\U0001F4E7  tanyasingh9270@gmail.com\n"
            "\U0001F310  github.com/Tanya9270\n"
            "\U0001F4BC  linkedin.com/in/tanya-singh",
            font_size=14, color=OFF_WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 13)


# ── Save ─────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "QR_Attendance_System_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
