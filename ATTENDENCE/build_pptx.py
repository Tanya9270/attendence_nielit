"""
Populate Dissertation_template.pptx with research content from SIUS paper.
Preserves all template formatting; only replaces placeholder text.
"""

from pptx import Presentation
from pptx.util import Pt
import copy

prs = Presentation('Dissertation_template.pptx')
slides = list(prs.slides)

# ─── Helper: clear a text-frame and set new text preserving first-run formatting ──
def set_text(shape, lines, font_size=None):
    """Replace all text in shape with `lines` (list of strings).
    Keeps the formatting of the first run of the first paragraph as template."""
    tf = shape.text_frame
    # Grab reference formatting from first paragraph / first run
    ref_para = tf.paragraphs[0]
    ref_run = ref_para.runs[0] if ref_para.runs else None

    # Clear existing paragraphs (keep the first, remove extras)
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
            # Clear existing runs
            for r in list(para.runs):
                r._r.getparent().remove(r._r)
            run = para.add_run()
        else:
            para = tf.add_paragraph()
            # Copy paragraph-level formatting from reference
            if ref_para.alignment is not None:
                para.alignment = ref_para.alignment
            para.level = ref_para.level
            run = para.add_run()

        run.text = line
        # Copy font formatting from reference run
        if ref_run is not None:
            if ref_run.font.name:
                run.font.name = ref_run.font.name
            if ref_run.font.bold is not None:
                run.font.bold = ref_run.font.bold
            if ref_run.font.color and ref_run.font.color.rgb:
                run.font.color.rgb = ref_run.font.color.rgb
        if font_size:
            run.font.size = Pt(font_size)
        elif ref_run and ref_run.font.size:
            run.font.size = ref_run.font.size


def set_text_simple(shape, lines, font_size=None):
    """Simpler version: just set text, try to keep format."""
    tf = shape.text_frame
    # Clear
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
            for r in list(para.runs):
                r._r.getparent().remove(r._r)
            run = para.add_run()
        else:
            para = tf.add_paragraph()
            run = para.add_run()
        run.text = line
        if font_size:
            run.font.size = Pt(font_size)


def get_shape_by_name(slide, name_fragment):
    """Find shape whose name contains the fragment."""
    for s in slide.shapes:
        if name_fragment in s.name:
            return s
    return None


def get_text_shapes(slide):
    """Get all shapes with text frames, in order."""
    return [s for s in slide.shapes if s.has_text_frame]


# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════
s = slides[0]
txt_shapes = get_text_shapes(s)
# Shape 0 (;89): Title
set_text_simple(txt_shapes[0], ["Prediction and Prevention of Diabetes using Explainable AI"])
# Shape 1 (;90): Subtitle
set_text_simple(txt_shapes[1], ["Dissertation Presentation"])
# Shape 2 (;91): Supervisor info
set_text_simple(txt_shapes[2], [
    "Under Supervision of:",
    "Dr. Rakesh Kumar",
    "Assistant Professor, Dept. of CSE"
])
# Shape 3 (;92): Student info
set_text_simple(txt_shapes[3], [
    "Tanya Singh",
    "Roll No: 215/ICS/030",
    "Integrated B.Tech. + M.Tech (CSE - Data Science)"
])
# Shape 5 (;94): slide number
set_text_simple(txt_shapes[5], ["1/21"])
# Shape 6 (;95): Date
set_text_simple(txt_shapes[6], ["February 2026"])

# ═══════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ═══════════════════════════════════════════════════════
s = slides[1]
txt_shapes = get_text_shapes(s)
# Title stays "Outline of the Presentation"
# Content shape (;105)
set_text_simple(txt_shapes[1], [
    "1. Introduction — Research Area & Problem Statement",
    "2. Background — Literature Review, Research Gaps & Objectives",
    "3. Methodology — Dataset, Feature Engineering, Ensemble Architecture",
    "4. Results & Discussion — Multi-Threshold Performance, Cross-Validation",
    "5. Conclusion — Key Findings, Limitations & Future Work",
    "6. References"
], font_size=14)
# Slide number
for sh in txt_shapes:
    if sh.text.strip() == "2":
        set_text_simple(sh, ["2/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 3 — Introduction Section Header
# ═══════════════════════════════════════════════════════
s = slides[2]
txt_shapes = get_text_shapes(s)
# Title stays "Introduction"
# Subtitle shape
set_text_simple(txt_shapes[1], [
    "Research Area",
    "Problem Statement"
])
for sh in txt_shapes:
    if sh.text.strip() == "3":
        set_text_simple(sh, ["3/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 4 — Research Area
# ═══════════════════════════════════════════════════════
s = slides[3]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "• Type 2 Diabetes Mellitus (T2DM) affects 537 million adults globally, projected to reach 783 million by 2045",
    "• Early identification through non-invasive screening is critical yet underserved",
    "• AI and Machine Learning have shown efficacy for diabetes risk assessment using epidemiological datasets",
    "• Gradient-boosted methods (XGBoost, LightGBM, CatBoost) consistently outperform traditional classifiers",
    "• Explainable AI (XAI) is essential for clinical trust — SHAP, LIME, counterfactual explanations",
    "• Key challenge: bridging the gap between research prototypes and deployable clinical systems"
], font_size=14)
for sh in txt_shapes:
    if sh.text.strip() == "4":
        set_text_simple(sh, ["4/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 5 — Problem Statement
# ═══════════════════════════════════════════════════════
s = slides[4]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "1. Most approaches use a single binary classification threshold (HbA1c ≥ 5.7%), ignoring clinical heterogeneity of screening scenarios",
    "",
    "2. Few systems combine per-patient SHAP explanations with actionable counterfactual recommendations for lifestyle changes",
    "",
    "3. Persistent translation gap between research prototypes and production-grade clinical systems — lacking fairness monitoring, deployment architecture",
    "",
    "4. No empirical study on the impact of target threshold definition on model performance has been reported"
], font_size=13)
for sh in txt_shapes:
    if sh.text.strip() == "5":
        set_text_simple(sh, ["5/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 6 — Background Section Header
# ═══════════════════════════════════════════════════════
s = slides[5]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Literature Review",
    "Research Gaps",
    "Research Objectives"
])
for sh in txt_shapes:
    if sh.text.strip() == "6":
        set_text_simple(sh, ["6/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 7 — Literature Review
# ═══════════════════════════════════════════════════════
s = slides[6]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "• Khokhar et al. (2025): Systematic review of 53 studies — gradient-boosted methods outperform traditional classifiers",
    "• Nie et al. (2024) — DRPM: Deep learning on NHANES, 36→5 features, strong AUC but limited feature set",
    "• Noh & Kim (2024): Gradient boosting + causal discovery (LiNGAM), ~84.8% accuracy",
    "• Shah et al. (2025): Hybrid ensemble for cardiovascular risk, AUC = 0.82, SHAP-based",
    "• Tasin et al. (2023): SHAP + LIME for diabetes prediction, improved clinical trust",
    "• Maimaitijiang et al. (2025): CatBoost + SHAP + personalised chatbot assistant",
    "• Kabir et al. (2025): Proximity-constrained counterfactuals for diabetes risk",
    "• Mohanty et al. (2025): IMPACT — counterfactuals + multimodal LLM for multi-disease prevention"
], font_size=12)
for sh in txt_shapes:
    if sh.text.strip() == "7":
        set_text_simple(sh, ["7/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 8 — Research Gaps
# ═══════════════════════════════════════════════════════
s = slides[7]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Gap 1: Single-threshold classification — majority of systems use only HbA1c ≥ 5.7% without multi-threshold adaptability",
    "",
    "Gap 2: Lack of integrated explainability — few combine per-patient SHAP explanations with actionable counterfactual interventions",
    "",
    "Gap 3: Research-to-deployment gap — most systems lack production-grade architecture, fairness monitoring, or HIPAA compliance",
    "",
    "Gap 4: No empirical study on how target threshold definition impacts predictive performance across identical architectures"
], font_size=13)
for sh in txt_shapes:
    if sh.text.strip() == "8":
        set_text_simple(sh, ["8/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 9 — Research Objectives
# ═══════════════════════════════════════════════════════
s = slides[8]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "1. Multi-Threshold Ensemble Architecture: Train three independent weighted-blend ensembles (LightGBM, CatBoost, XGBoost) for HbA1c thresholds ≥5.7%, ≥5.9%, ≥6.5%",
    "",
    "2. Target Threshold Dominance Discovery: Empirically demonstrate that threshold choice is the single most impactful determinant of AUC (0.885 → 0.972)",
    "",
    "3. Integrated Multi-Modal Explainability: SHAP waterfall + DiCE counterfactual interventions + WHO epidemiological cross-validation",
    "",
    "4. Production-Grade Deployment: Full-stack system — FastAPI backend, Next.js 14 frontend, fairness auditing, HIPAA-compliant, containerised"
], font_size=13)
for sh in txt_shapes:
    if sh.text.strip() == "9":
        set_text_simple(sh, ["9/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 10 — Methodology Section Header
# ═══════════════════════════════════════════════════════
s = slides[9]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Dataset & Data Preparation",
    "Feature Engineering & Multi-Threshold Design",
    "Ensemble Architecture & Explainability"
])
for sh in txt_shapes:
    if sh.text.strip() == "10":
        set_text_simple(sh, ["10/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 11 — Dataset
# ═══════════════════════════════════════════════════════
s = slides[10]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "• Source: U.S. National Health and Nutrition Examination Survey (NHANES)",
    "• Span: 12 survey cycles (1999–2024), 6 domains merged on SEQN",
    "• Initial pool: 67,876 records → Filtered to 33,333 (fasting glucose subsample)",
    "• Filtering rationale: 50.9% fasting glucose missing; median imputation would cap AUC at ~0.86",
    "• Split: 80/20 stratified on binary target (26,666 train / 6,667 test)",
    "• Optuna HPO: further 60/20/20 train/val/test split",
    "• 62 total features: 23 base + 23 engineered interactions + 16 missingness indicators",
    "• HbA1c excluded from features (used as target) — FastingGlucose serves as primary glycaemic predictor"
], font_size=13)
for sh in txt_shapes:
    if sh.text.strip() == "11":
        set_text_simple(sh, ["11/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 12 — Methodology
# ═══════════════════════════════════════════════════════
s = slides[11]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Multi-Threshold Target Definition:",
    "  • ≥5.7% (ADA Prediabetes, 35.5% prevalence)",
    "  • ≥5.9% (High-Risk Prediabetes, 24.2% prevalence) — production default",
    "  • ≥6.5% (Diabetes Diagnosis, 11.0% prevalence)",
    "",
    "Ensemble Architecture:",
    "  • Weighted blend: p̂(x) = w₁·LGB(x) + w₂·CB(x) + w₃·XGB(x)",
    "  • Weights optimised via grid search (w₁+w₂+w₃ = 1)",
    "  • 5× faster than stacking, simpler deployment, more interpretable",
    "",
    "Key Design Decisions:",
    "  • Median imputation > MICE (avoids synthetic inter-feature correlations)",
    "  • Clinical range clipping: BMI ∈ [12,70], FastingGlucose ∈ [40,500] mg/dL"
], font_size=12)
for sh in txt_shapes:
    if sh.text.strip() == "12":
        set_text_simple(sh, ["12/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 13 — Model Architecture/Diagram
# ═══════════════════════════════════════════════════════
s = slides[12]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Three-Tier Explainability Framework:",
    "",
    "Tier 1 — Individual-Level: TreeExplainer computes exact Shapley values → per-patient waterfall plots on 12 key features",
    "",
    "Tier 2 — Actionable Counterfactuals: DiCE generates minimal modifiable feature changes (BMI↓, Fiber↑, Sleep→7-8h) with SHAP-guided greedy search",
    "",
    "Tier 3 — Population-Level: WHO epidemiological cross-validation verifies model feature importance aligns with global diabetes prevalence patterns",
    "",
    "Deployment: FastAPI + Uvicorn backend, Next.js 14 frontend, PostgreSQL 18, Docker Compose (6 services), Prometheus + Grafana monitoring"
], font_size=12)
for sh in txt_shapes:
    if sh.text.strip() == "13":
        set_text_simple(sh, ["13/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 14 — Results & Discussion Section Header
# ═══════════════════════════════════════════════════════
s = slides[13]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Multi-Threshold Performance",
    "Cross-Validation Stability",
    "Model Evolution & Comparison"
])
for sh in txt_shapes:
    if sh.text.strip() == "14":
        set_text_simple(sh, ["14/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 15 — Results & Discussion
# ═══════════════════════════════════════════════════════
s = slides[14]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Multi-Threshold Production Performance:",
    "  • ≥5.7%: Blend AUC = 0.885 | CV Mean = 0.891 | Prevalence = 35.5%",
    "  • ≥5.9%: Blend AUC = 0.923 | CV Mean = 0.923 | Prevalence = 24.2%  ← Production Default",
    "  • ≥6.5%: Blend AUC = 0.972 | Prevalence = 11.0%",
    "",
    "Key Finding — Threshold Dominance:",
    "  • AUC increases monotonically with threshold: class separability > model complexity",
    "  • CatBoost dominates at ≥5.9% (w=0.45), LightGBM at ≥5.7% (w=0.70)",
    "",
    "5-Fold CV Stability: Mean AUC = 0.889 ± 0.005, Sensitivity = 0.933 ± 0.010, Brier = 0.134 ± 0.003",
    "",
    "Model evolved through 7 rounds: 0.812 → 0.923 AUC; largest gain from threshold sweep (+0.029), not architecture"
], font_size=12)
for sh in txt_shapes:
    if sh.text.strip() == "15":
        set_text_simple(sh, ["15/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 16 — Conclusion Section Header
# ═══════════════════════════════════════════════════════
s = slides[15]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Conclusion and Future Scope",
    "List of Publications",
    "Next Semester Plan"
])
for sh in txt_shapes:
    if sh.text.strip() == "16":
        set_text_simple(sh, ["16/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 17 — Conclusion and Future Scope
# ═══════════════════════════════════════════════════════
s = slides[16]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "Key Contributions:",
    "  1. Multi-threshold design: 3 HbA1c thresholds (≥5.7%, ≥5.9%, ≥6.5%) — AUCs of 0.885, 0.923, 0.972",
    "  2. Target threshold dominates model architecture in determining performance",
    "  3. Three-tier explainability: SHAP + DiCE counterfactuals + WHO cross-validation",
    "  4. Production-grade deployment with fairness monitoring & HIPAA compliance",
    "",
    "Limitations: NHANES-only (U.S. population), AUC bounded by non-HbA1c feature content, no external clinical validation",
    "",
    "Future Work: External validation on non-U.S. populations, longitudinal risk modelling, agentic AI integration, CGM integration, federated learning, causal counterfactuals"
], font_size=12)
for sh in txt_shapes:
    if sh.text.strip() == "17":
        set_text_simple(sh, ["17/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 18 — List of Publications
# ═══════════════════════════════════════════════════════
s = slides[17]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "T. Singh and R. Kumar, \"Prediction and Prevention of Diabetes using Explainable AI,\" submitted to IEEE conference proceedings, 2026."
], font_size=14)
for sh in txt_shapes:
    if sh.text.strip() == "18":
        set_text_simple(sh, ["18/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 19 — Next Semester Plan
# ═══════════════════════════════════════════════════════
s = slides[18]
txt_shapes = get_text_shapes(s)
set_text_simple(txt_shapes[1], [
    "• External validation on international clinical datasets (non-U.S. populations)",
    "• Integration of longitudinal NHANES data for temporal risk trajectory modelling",
    "• Agentic AI module for automated follow-up scheduling and adaptive intervention planning",
    "• Continuous Glucose Monitoring (CGM) integration using transformer-based deep learning",
    "• Expand counterfactual engine with causal discovery methods (LiNGAM)",
    "• Federated learning for multi-institutional training without data centralisation",
    "• Personalised threshold selection algorithms based on patient clinical profile"
], font_size=13)
for sh in txt_shapes:
    if sh.text.strip() == "19":
        set_text_simple(sh, ["19/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 20 — References
# ═══════════════════════════════════════════════════════
s = slides[19]
txt_shapes = get_text_shapes(s)
# The second text shape (;393) is the content area
for sh in txt_shapes:
    if '393' in sh.name:
        set_text_simple(sh, [
            "[1] Khokhar et al., \"Advances in AI for diabetes prediction,\" Artif. Intell. Med., 2025.",
            "[2] Khalifa & Albadawy, \"AI for diabetes: prevention, diagnosis, management,\" Comput. Methods Programs Biomed., 2024.",
            "[3] Nie et al., \"DRPM: Advanced predictive model for early diabetes detection,\" Heliyon, 2024.",
            "[4] Bontha et al., \"Predicting risk and complications of diabetes through AI,\" Appl. Sci., 2024.",
            "[5] Tasin et al., \"Diabetes prediction using ML and explainable AI,\" Healthcare Tech. Lett., 2023.",
            "[6] Noh & Kim, \"Diabetes prediction through causal discovery and ML,\" Appl. Sci., 2024.",
            "[7] Shah et al., \"Cardiovascular risk with hybrid ensemble and XAI,\" Sci. Rep., 2025.",
            "[8] Prakash & Subhajini, \"Diabetics prediction with LwCPSA and XAI,\" IJPRAI, 2025.",
            "[9] Maimaitijiang et al., \"XAI framework for diabetes with chatbot,\" Healthcare Analytics, 2025.",
            "[10] Kabir et al., \"Proximity-constrained counterfactuals for diabetes,\" ACISC, 2025.",
            "[11] Mohanty et al., \"IMPACT: Multi-disease prevention with XAI and LLM,\" PeerJ CS, 2025.",
            "[12] Prashanthan & Prashanthan, \"T2DM risk in gestational diabetes history,\" Primary Care Diabetes, 2025."
        ], font_size=11)
        break
# Slide number
for sh in txt_shapes:
    if sh.text.strip() == "20":
        set_text_simple(sh, ["20/21"])
        break

# ═══════════════════════════════════════════════════════
# SLIDE 21 — Thank You (keep as-is, just update slide number)
# ═══════════════════════════════════════════════════════
s = slides[20]
txt_shapes = get_text_shapes(s)
for sh in txt_shapes:
    if sh.text.strip() == "21":
        set_text_simple(sh, ["21/21"])
        break

# ─── Save ──────────────────────────────────────────────
prs.save('Dissertation_Presentation.pptx')
print("✅ Saved: Dissertation_Presentation.pptx")
