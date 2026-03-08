from pptx import Presentation
from pptx.util import Emu
from lxml import etree

prs = Presentation('Dissertation_template.pptx')

# Check background of different slides
for i, slide in enumerate(prs.slides, 1):
    bg = slide.background
    fill = bg.fill
    try:
        print(f'Slide {i} bg fill type: {fill.type}')
        if fill.type is not None:
            try:
                print(f'  fore color: {fill.fore_color.rgb}')
            except:
                print(f'  fore color: N/A')
    except:
        print(f'Slide {i} bg: default')

# Get the header bar (group shape) details from slide 2
s2 = list(prs.slides)[1]
for sh in s2.shapes:
    if sh.shape_type == 6:  # GROUP
        print(f'\nHeader group: pos=({sh.left},{sh.top}), size=({sh.width},{sh.height})')
        grp = sh._element
        for child in grp:
            tag = etree.QName(child.tag).localname
            if tag == 'sp':
                for desc in child.iter():
                    qtag = etree.QName(desc.tag).localname
                    if qtag == 'srgbClr':
                        print(f'  Group child srgb color: {desc.get("val")}')
                    elif qtag == 'schemeClr':
                        print(f'  Group child scheme color: {desc.get("val")}')
            elif tag == 'pic':
                print(f'  Group child: picture (logo)')

# Check section header slide (slide 3) for styling
print('\n--- Section header slide styling ---')
s3 = list(prs.slides)[2]
for sh in s3.shapes:
    if sh.has_text_frame:
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                try:
                    tc = run.font.color.theme_color
                    print(f'Slide 3 text "{run.text[:30]}" theme_color={tc}')
                except:
                    pass
                try:
                    print(f'  rgb={run.font.color.rgb}')
                except:
                    pass

# Dig into the XML for theme colors
print('\n--- Theme colors from slide master ---')
master = prs.slide_masters[0]
theme_elem = master.element
ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
for clr in theme_elem.iter():
    qtag = etree.QName(clr.tag).localname
    if qtag in ('dk1', 'dk2', 'lt1', 'lt2', 'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6', 'hlink', 'folHlink'):
        for child in clr:
            ctag = etree.QName(child.tag).localname
            if ctag == 'srgbClr':
                print(f'  {qtag}: #{child.get("val")}')
            elif ctag == 'sysClr':
                print(f'  {qtag}: system {child.get("lastClr")}')

# Check subtitle color on slide 1
print('\n--- Slide 1 subtitle color ---')
s1 = list(prs.slides)[0]
for sh in s1.shapes:
    if sh.has_text_frame:
        txt = sh.text_frame.text[:40]
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                try:
                    rgb = run.font.color.rgb
                    print(f'  "{txt}" run color: {rgb}')
                except:
                    pass
                try:
                    tc = run.font.color.theme_color
                    print(f'  "{txt}" theme color: {tc}')
                except:
                    pass
